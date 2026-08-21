"""Text extraction from PDF and PPTX files.

This module is the parsing agent's core: it turns a raw submission file
(PDF or PPTX) into structured, queryable text. PDFs are extracted page by
page with PyMuPDF; PPTX decks are extracted slide by slide with python-pptx
(including titles, body text, and speaker notes).
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from typing import Any

import fitz  # PyMuPDF
from pptx import Presentation


class UnsupportedFormatError(ValueError):
    """Raised when the file type is not PDF or PPTX."""


class ParsingError(RuntimeError):
    """Raised when a file cannot be parsed (corrupt, empty, etc.)."""


@dataclass
class ParsedDocument:
    """Structured output of the parsing agent.

    Attributes:
        source_format: One of ``pdf`` or ``pptx``.
        raw_text: The full extracted text joined together.
        sections: A list of per-page (PDF) or per-slide (PPTX) text chunks.
    """

    source_format: str
    raw_text: str
    sections: list[dict[str, Any]] = field(default_factory=list)


def _extract_pdf(file_bytes: bytes) -> ParsedDocument:
    """Extract text from a PDF, one section per page."""
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        raise ParsingError(f"Could not open PDF: {exc}") from exc

    sections: list[dict[str, Any]] = []
    page_texts: list[str] = []

    try:
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            page_texts.append(text)
            sections.append({"type": "page", "index": page_num, "text": text})
    finally:
        doc.close()

    raw_text = "\n\n".join(page_texts).strip()
    return ParsedDocument(source_format="pdf", raw_text=raw_text, sections=sections)


def _extract_pptx(file_bytes: bytes) -> ParsedDocument:
    """Extract text from a PPTX deck, one section per slide.

    Includes slide titles, body text, and speaker notes.
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:
        raise ParsingError(f"Could not open PPTX: {exc}") from exc

    sections: list[dict[str, Any]] = []
    slide_texts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        # Title (if present).
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            parts.append(slide.shapes.title.text.strip())

        # All other text-bearing shapes (text boxes, body placeholders, etc.).
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

        # Speaker notes.
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Notes] {notes_text}")

        slide_text = "\n".join(parts).strip()
        slide_texts.append(slide_text)
        sections.append({"type": "slide", "index": slide_num, "text": slide_text})

    raw_text = "\n\n".join(slide_texts).strip()
    return ParsedDocument(source_format="pptx", raw_text=raw_text, sections=sections)


def extract_text(file_bytes: bytes, source_format: str) -> ParsedDocument:
    """Extract structured text from a PDF or PPTX file.

    Args:
        file_bytes: The raw file contents.
        source_format: One of ``pdf`` or ``pptx``.

    Returns:
        A :class:`ParsedDocument` with ``raw_text`` and per-page/slide sections.

    Raises:
        UnsupportedFormatError: If ``source_format`` is not pdf/pptx.
        ParsingError: If the file cannot be parsed.
    """
    if source_format == "pdf":
        return _extract_pdf(file_bytes)
    if source_format == "pptx":
        return _extract_pptx(file_bytes)
    raise UnsupportedFormatError(
        f"Unsupported format '{source_format}'. Only 'pdf' and 'pptx' are supported."
    )
