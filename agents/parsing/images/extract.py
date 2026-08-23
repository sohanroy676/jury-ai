"""Embedded-image extraction from PDF and PPTX files (v0.3.5).

Extracts images placed in page/slide CONTENT areas. Applies the
roadmap's structural filter for PPTX: only shapes on the slide itself
are considered — images living in the slide master or layout (logos,
backgrounds, theme graphics) are never visited, so they can never
reach the classifier or the vision model. PDFs have no master/layout
concept, so every embedded image is a candidate.

Images smaller than ``min_dimension`` on either side are treated as
icons/bullets and skipped.
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass

import pymupdf  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from agents.parsing.extractor import UnsupportedFormatError

logger = logging.getLogger(__name__)

# Images smaller than this on either side are treated as icons/bullets
# and skipped. Overridable via the MIN_IMAGE_DIMENSION env var at the
# pipeline/route level; this default keeps extraction self-contained.
DEFAULT_MIN_DIMENSION = 64


@dataclass
class ExtractedImage:
    """One embedded image candidate from a submission file."""

    page_number: int  # 1-based page (PDF) or slide (PPTX) index
    image_bytes: bytes
    source_format: str  # "pdf" or "pptx"


def _extract_pdf_images(file_bytes: bytes, min_dimension: int) -> list[ExtractedImage]:
    """Extract embedded images from a PDF, one candidate per unique xref."""
    try:
        doc = pymupdf.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        # Let the text extractor own PDF-open error semantics; image
        # extraction is only ever called after text parsing succeeded.
        raise ValueError(f"Could not open PDF for image extraction: {exc}") from exc

    results: list[ExtractedImage] = []

    try:
        for page_num, page in enumerate(doc, start=1):
            seen_xrefs: set[int] = set()
            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)

                try:
                    info = doc.extract_image(xref)
                except (RuntimeError, ValueError) as exc:
                    # Unreadable/corrupt embedded image — skip it rather
                    # than failing the whole document.
                    logger.warning(
                        "Skipping unreadable embedded image (xref %s): %s",
                        xref,
                        exc,
                    )
                    continue

                width = int(info.get("width") or 0)
                height = int(info.get("height") or 0)
                if width < min_dimension or height < min_dimension:
                    continue

                results.append(
                    ExtractedImage(
                        page_number=page_num,
                        image_bytes=info["image"],
                        source_format="pdf",
                    )
                )
    finally:
        doc.close()

    return results


def _iter_picture_shapes(shapes):
    """Yield picture shapes, recursing into group shapes."""
    for shape in shapes:
        try:
            shape_type = shape.shape_type
        except (AttributeError, TypeError, ValueError) as exc:
            # Exotic shapes can raise on shape_type — skip them.
            logger.warning("Skipping shape with unreadable type: %s", exc)
            continue
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _extract_pptx_images(file_bytes: bytes, min_dimension: int) -> list[ExtractedImage]:
    """Extract pictures from PPTX slides (content areas only).

    Structural filter: iterating ``slide.shapes`` never sees shapes on
    the slide master or layout, so template logos/backgrounds are
    excluded by construction.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    results: list[ExtractedImage] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in _iter_picture_shapes(slide.shapes):
            try:
                image = shape.image
                image_bytes = image.blob
                width, height = image.size
            except (OSError, ValueError) as exc:
                # Unsupported picture format (e.g. some EMF/SVG edges)
                # or unreadable blob — skip rather than fail the deck.
                # python-pptx 1.x delegates decoding to Pillow, so
                # unknown formats surface as UnidentifiedImageError
                # (an OSError subclass).
                logger.warning(
                    "Skipping unreadable slide picture on slide %s: %s",
                    slide_num,
                    exc,
                )
                continue

            if width < min_dimension or height < min_dimension:
                continue

            results.append(
                ExtractedImage(
                    page_number=slide_num,
                    image_bytes=image_bytes,
                    source_format="pptx",
                )
            )

    return results


def extract_images(
    file_bytes: bytes,
    source_format: str,
    min_dimension: int = DEFAULT_MIN_DIMENSION,
) -> list[ExtractedImage]:
    """Extract embedded content-area images from a PDF or PPTX file.

    Args:
        file_bytes: The raw file contents.
        source_format: One of ``pdf`` or ``pptx``.
        min_dimension: Skip images smaller than this on either side.

    Returns:
        A list of :class:`ExtractedImage` candidates (may be empty).

    Raises:
        UnsupportedFormatError: If ``source_format`` is not pdf/pptx.
    """
    if source_format == "pdf":
        return _extract_pdf_images(file_bytes, min_dimension)
    if source_format == "pptx":
        return _extract_pptx_images(file_bytes, min_dimension)
    raise UnsupportedFormatError(
        f"Unsupported format '{source_format}'. Only 'pdf' and 'pptx' are supported."
    )
