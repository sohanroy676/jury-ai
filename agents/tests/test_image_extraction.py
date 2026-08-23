"""Tests for embedded-image extraction from PDF and PPTX (v0.3.5)."""

import copy
import io

import pymupdf
import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches

from agents.parsing.images.extract import ExtractedImage, extract_images

# --- Fixture helpers ---------------------------------------------------------


def _make_png(width: int = 200, height: int = 120, color=(180, 40, 40)) -> bytes:
    """Build a real in-memory PNG."""
    img = Image.new("RGB", (width, height), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_pdf_with_image(image_bytes: bytes, text: str | None = None) -> bytes:
    """Build a real in-memory PDF containing one embedded image."""
    doc = pymupdf.open()
    page = doc.new_page()
    if text:
        page.insert_text((72, 72), text)
    page.insert_image(pymupdf.Rect(72, 100, 272, 220), stream=image_bytes)
    data = doc.tobytes()
    doc.close()
    return data


def _make_pptx_with_picture(image_bytes: bytes) -> bytes:
    """Build a real in-memory PPTX with one picture on slide 1."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(1), Inches(1))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_master_picture(image_bytes: bytes) -> bytes:
    """Build a PPTX whose picture lives on the SLIDE MASTER (template logo).

    This is the structural-filter fixture: the image must never be
    extracted because it is not part of any slide's content shapes.

    MasterShapes has no add_picture API, so the picture is added to a
    slide first, then its XML element is reparented into the master's
    shape tree with a fresh master-part relationship for the image.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(1), Inches(1))

    master = prs.slide_masters[0]
    image_part = slide.part.related_part(pic._element.blip_rId)
    new_rid = master.part.relate_to(image_part, RT.IMAGE)

    pic_copy = copy.deepcopy(pic._element)
    pic_copy.blipFill.blip.rEmbed = new_rid
    master.shapes._spTree.append(pic_copy)

    # Remove the original from the slide so ONLY the master has it.
    pic._element.getparent().remove(pic._element)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_grouped_picture(image_bytes: bytes) -> bytes:
    """Build a PPTX with a picture nested inside a group shape."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.shapes.add_picture(io.BytesIO(image_bytes), Inches(1), Inches(1))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- PDF extraction ----------------------------------------------------------


def test_pdf_extracts_embedded_image():
    """An embedded PNG is extracted with the correct page number."""
    png = _make_png()
    data = _make_pdf_with_image(png, text="Architecture overview")

    results = extract_images(data, "pdf")

    assert len(results) == 1
    assert isinstance(results[0], ExtractedImage)
    assert results[0].page_number == 1
    assert results[0].source_format == "pdf"
    assert len(results[0].image_bytes) > 0


def test_pdf_without_images_returns_empty():
    """A text-only PDF yields no image candidates (no crash)."""
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Just text, no images here")
    data = doc.tobytes()
    doc.close()

    assert extract_images(data, "pdf") == []


def test_pdf_skips_tiny_images():
    """Images below the minimum dimension are treated as icons and skipped."""
    tiny = _make_png(width=20, height=10)
    data = _make_pdf_with_image(tiny)

    assert extract_images(data, "pdf") == []


def test_pdf_min_dimension_override_keeps_small_images():
    """Lowering min_dimension keeps images the default would skip."""
    tiny = _make_png(width=20, height=10)
    data = _make_pdf_with_image(tiny)

    results = extract_images(data, "pdf", min_dimension=10)
    assert len(results) == 1


def test_pdf_same_xreferenced_twice_extracted_once(monkeypatch):
    """The same embedded image referenced repeatedly is returned once."""
    data = _make_pdf_with_image(_make_png())

    original_get_images = pymupdf.Page.get_images

    def doubled_get_images(self, *args, **kwargs):
        images = list(original_get_images(self, *args, **kwargs))
        return images + images  # report every xref twice

    monkeypatch.setattr(pymupdf.Page, "get_images", doubled_get_images)

    results = extract_images(data, "pdf")
    assert len(results) == 1


def test_pdf_unreadable_embedded_image_skipped(monkeypatch):
    """A corrupt/unreadable embedded image is skipped, not fatal."""

    def broken_extract_image(self, xref):
        raise RuntimeError("simulated unreadable image")

    data = _make_pdf_with_image(_make_png())
    monkeypatch.setattr(pymupdf.Document, "extract_image", broken_extract_image)

    assert extract_images(data, "pdf") == []


def test_pdf_multi_page_reports_correct_pages():
    """Images on different pages carry their own page numbers."""
    png_a = _make_png(color=(200, 30, 30))
    png_b = _make_png(color=(30, 200, 30))

    doc = pymupdf.open()
    page1 = doc.new_page()
    page1.insert_text((72, 72), "Page one")
    page1.insert_image(pymupdf.Rect(72, 100, 272, 220), stream=png_a)
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Page two")
    page2.insert_image(pymupdf.Rect(72, 100, 272, 220), stream=png_b)
    data = doc.tobytes()
    doc.close()

    results = extract_images(data, "pdf")
    assert [r.page_number for r in results] == [1, 2]


# --- PPTX extraction ---------------------------------------------------------


def test_pptx_extracts_slide_picture():
    """A picture placed on a slide is extracted with its slide number."""
    data = _make_pptx_with_picture(_make_png())

    results = extract_images(data, "pptx")

    assert len(results) == 1
    assert results[0].page_number == 1
    assert results[0].source_format == "pptx"
    assert len(results[0].image_bytes) > 0


def test_pptx_master_picture_not_extracted():
    """Structural filter: slide-master/template logos are NEVER extracted."""
    data = _make_pptx_with_master_picture(_make_png())

    # Sanity: the fixture really does place a picture on the master.
    prs_check = Presentation(io.BytesIO(data))
    master_has_picture = any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        for shape in prs_check.slide_masters[0].shapes
    )
    assert master_has_picture, "fixture failed to place picture on master"

    # ...and the extractor must not return it.
    assert extract_images(data, "pptx") == []


def test_pptx_grouped_picture_extracted():
    """Pictures nested inside group shapes are still found."""
    data = _make_pptx_with_grouped_picture(_make_png())

    results = extract_images(data, "pptx")
    assert len(results) == 1


def test_pptx_without_pictures_returns_empty():
    """A text-only deck yields no image candidates (no crash)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Text only"
    slide.placeholders[1].text = "No pictures anywhere"
    buf = io.BytesIO()
    prs.save(buf)

    assert extract_images(buf.getvalue(), "pptx") == []


def test_pptx_multi_slide_reports_correct_slides():
    """Pictures on different slides carry their own slide numbers."""
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.shapes.add_picture(
        io.BytesIO(_make_png(color=(200, 30, 30))), Inches(1), Inches(1)
    )
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.shapes.add_picture(
        io.BytesIO(_make_png(color=(30, 30, 200))), Inches(1), Inches(1)
    )
    buf = io.BytesIO()
    prs.save(buf)

    results = extract_images(buf.getvalue(), "pptx")
    assert [r.page_number for r in results] == [1, 2]


# --- Format handling ---------------------------------------------------------


def test_unsupported_format_raises():
    """An unknown source_format raises UnsupportedFormatError."""
    from agents.parsing.extractor import UnsupportedFormatError

    with pytest.raises(UnsupportedFormatError):
        extract_images(b"whatever", "docx")
