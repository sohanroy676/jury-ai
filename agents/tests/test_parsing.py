"""Tests for the parsing agent (PDF + PPTX text extraction)."""

import io

import fitz
import pytest
from pptx import Presentation

from agents.parsing.extractor import (
    ParsingError,
    UnsupportedFormatError,
    extract_text,
)


def _make_pdf(texts: list[str]) -> bytes:
    """Build a real in-memory PDF with one page per string."""
    doc = fitz.open()
    for text in texts:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    data = doc.tobytes()
    doc.close()
    return data


def _make_pptx(slides: list[dict]) -> bytes:
    """Build a real in-memory PPTX deck.

    Each dict may contain ``title``, ``body`` (list of strings), and ``notes``.
    """
    prs = Presentation()
    for spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        if spec.get("title"):
            slide.shapes.title.text = spec["title"]
        for i, line in enumerate(spec.get("body", [])):
            # Use the content placeholder's text frame.
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            if i == 0:
                tf.text = line
            else:
                tf.add_paragraph().text = line
        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- PDF tests -------------------------------------------------------------


def test_pdf_extracts_text_per_page():
    """A multi-page PDF yields one section per page with correct text."""
    data = _make_pdf(["Page one content", "Page two content"])
    result = extract_text(data, "pdf")

    assert result.source_format == "pdf"
    assert "Page one content" in result.raw_text
    assert "Page two content" in result.raw_text
    assert len(result.sections) == 2
    assert result.sections[0]["type"] == "page"
    assert result.sections[0]["index"] == 1
    assert result.sections[1]["index"] == 2


def test_pdf_blank_page_yields_empty_raw_text():
    """A PDF with a blank page yields empty raw_text (no crash)."""
    doc = fitz.open()
    doc.new_page()  # A page with no text at all.
    data = doc.tobytes()
    doc.close()

    result = extract_text(data, "pdf")
    assert result.raw_text == ""
    assert len(result.sections) == 1
    assert result.sections[0]["text"] == ""


def test_pdf_corrupt_bytes_raises():
    """Garbage bytes that aren't a PDF raise ParsingError."""
    with pytest.raises(ParsingError):
        extract_text(b"this is not a pdf at all", "pdf")


def test_pdf_image_only_page_has_empty_text():
    """A page with no extractable text yields an empty section, not a crash."""
    doc = fitz.open()
    page = doc.new_page()
    # Draw a rectangle (no text) so the page exists but has no text.
    page.draw_rect(fitz.Rect(72, 72, 200, 200), color=(0, 0, 0))
    data = doc.tobytes()
    doc.close()

    result = extract_text(data, "pdf")
    assert len(result.sections) == 1
    assert result.sections[0]["text"] == ""


# --- PPTX tests ------------------------------------------------------------


def test_pptx_extracts_title_body_and_notes():
    """A deck with title, body, and notes extracts all three."""
    data = _make_pptx(
        [
            {
                "title": "Our Solution",
                "body": ["We built a platform", "It scales to millions"],
                "notes": "Demo: show the dashboard",
            }
        ]
    )
    result = extract_text(data, "pptx")

    assert result.source_format == "pptx"
    assert "Our Solution" in result.raw_text
    assert "We built a platform" in result.raw_text
    assert "It scales to millions" in result.raw_text
    assert "[Notes] Demo: show the dashboard" in result.raw_text
    assert len(result.sections) == 1
    assert result.sections[0]["type"] == "slide"
    assert result.sections[0]["index"] == 1


def test_pptx_multiple_slides():
    """A multi-slide deck yields one section per slide."""
    data = _make_pptx(
        [
            {"title": "Slide A", "body": ["Alpha"]},
            {"title": "Slide B", "body": ["Beta"]},
        ]
    )
    result = extract_text(data, "pptx")

    assert len(result.sections) == 2
    assert result.sections[0]["index"] == 1
    assert result.sections[1]["index"] == 2
    assert "Alpha" in result.sections[0]["text"]
    assert "Beta" in result.sections[1]["text"]


def test_pptx_slide_without_title():
    """A slide with no title still extracts body text."""
    data = _make_pptx([{"body": ["Just body text here"]}])
    result = extract_text(data, "pptx")

    assert "Just body text here" in result.raw_text
    assert len(result.sections) == 1


def test_pptx_corrupt_bytes_raises():
    """Garbage bytes that aren't a PPTX raise ParsingError."""
    with pytest.raises(ParsingError):
        extract_text(b"this is not a pptx at all", "pptx")


# --- Unicode dash normalization ---------------------------------------------


def test_normalize_unicode_dashes_all_variants():
    """All six typographic dash variants map to a plain ASCII hyphen."""
    from agents.parsing.extractor import normalize_unicode_dashes

    text = "a\u2010b\u2011c\u2012d\u2013e\u2014f\u2212g"
    assert normalize_unicode_dashes(text) == "a-b-c-d-e-f-g"
    assert normalize_unicode_dashes("no dashes here") == "no dashes here"


# NOTE: There is no PDF equivalent of the PPTX normalization test below.
# PyMuPDF's built-in fonts (helv, china-s) cannot encode typographic dash
# glyphs at write time — they are substituted with U+FFFD in the fixture,
# so a round-trip test cannot exercise them. The normalization itself is
# format-independent: the same shared function is applied to extracted
# text in both _extract_pdf and _extract_pptx, and is fully covered by
# test_normalize_unicode_dashes_all_variants plus this end-to-end test.


def test_pptx_normalizes_unicode_dashes():
    """Typographic dashes in a PPTX come out as plain ASCII hyphens."""
    data = _make_pptx(
        [
            {
                "title": "Cost\u2011per\u2011dollar",
                "body": [
                    "end\u2010to\u2010end pipeline",
                    "range \u2013 wide \u2014 scope",
                ],
            }
        ]
    )
    result = extract_text(data, "pptx")

    assert "Cost-per-dollar" in result.raw_text
    assert "end-to-end pipeline" in result.raw_text
    assert "range - wide - scope" in result.raw_text
    for ch in ("\u2010", "\u2011", "\u2012", "\u2013", "\u2014"):
        assert ch not in result.raw_text


# --- Format handling -------------------------------------------------------


def test_unsupported_format_raises():
    """An unknown source_format raises UnsupportedFormatError."""
    with pytest.raises(UnsupportedFormatError):
        extract_text(b"whatever", "docx")


def test_empty_bytes_raises_for_pdf():
    """Empty bytes cannot be parsed as a PDF."""
    with pytest.raises(ParsingError):
        extract_text(b"", "pdf")
