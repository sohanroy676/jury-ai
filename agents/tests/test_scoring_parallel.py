"""Orchestration tests for v0.5.0 scoring: parallel fan-out, ordering,
fail-closed semantics, and format independence (PPTX-only path)."""

from __future__ import annotations

import asyncio
import json
import time
from io import BytesIO
from unittest.mock import Mock

import pytest
from groq import RateLimitError
from pptx import Presentation
from pptx.util import Inches

from agents.parsing.extractor import extract_text
from agents.scoring.base import CRITERIA_NAMES, CriterionScore, ScoringResult
from agents.scoring.scorer import build_scoring_text, score_submission

MARKER = "QuantumQuokka pptx-only proposal"

DISTINCT_SCORES = {
    "problem_fit": 3,
    "technical_depth": 9,
    "feasibility": 6,
    "innovation": 8,
}


# --- Helpers ------------------------------------------------------------------


def _mock_response(content: str) -> Mock:
    resp = Mock()
    resp.choices = [Mock()]
    resp.choices[0].message.content = content
    return resp


def _valid_json(criterion: str, score: int | None = None) -> str:
    return json.dumps(
        {
            criterion: {
                "score": score if score is not None else DISTINCT_SCORES[criterion],
                "justification": f"{criterion} ok",
            }
        }
    )


@pytest.fixture(autouse=True)
def _no_backoff_sleep(monkeypatch):
    """Backoff waits never actually wait; concurrency tests sleep inside
    their fake clients directly."""

    async def _instant(_seconds: float) -> None:
        return None

    monkeypatch.setattr("agents.scoring.base._sleep", _instant)


@pytest.fixture(autouse=True)
def _clean_env(monkeypatch):
    monkeypatch.delenv("GROQ_API_KEY", raising=False)


def _router_client(on_create=None):
    """Fake AsyncGroq client that answers each specialist for its OWN
    criterion by inspecting the system prompt. ``on_create`` receives the
    criterion just before responding (used for concurrency tracking).
    """

    async def create(**kwargs):
        system_content = kwargs["messages"][0]["content"]
        criterion = next(c for c in CRITERIA_NAMES if c in system_content)
        if on_create is not None:
            on_create(criterion)
        return _mock_response(_valid_json(criterion))

    client = Mock()
    client.chat.completions.create.side_effect = create
    return client


def _client_failing_criterion(fail_criterion: str, bad_effect):
    """All specialists succeed except ``fail_criterion``, whose every call
    hits ``bad_effect`` (a raw-content string or Exception instance)."""

    async def create(**kwargs):
        system_content = kwargs["messages"][0]["content"]
        criterion = next(c for c in CRITERIA_NAMES if c in system_content)
        if criterion == fail_criterion:
            if isinstance(bad_effect, BaseException):
                raise bad_effect
            return _mock_response(bad_effect)
        return _mock_response(_valid_json(criterion))

    client = Mock()
    client.chat.completions.create.side_effect = create
    return client


# --- Aggregation & ordering ----------------------------------------------------


def test_all_four_criteria_aggregated_in_canonical_order(monkeypatch):
    """Result lists all four scores ordered canonically, whatever the
    internal completion order."""
    # technical_depth answers instantly; others lag behind it.
    state = {"active": 0, "max_active": 0}

    async def create(**kwargs):
        system_content = kwargs["messages"][0]["content"]
        criterion = next(c for c in CRITERIA_NAMES if c in system_content)
        state["active"] += 1
        state["max_active"] = max(state["max_active"], state["active"])
        try:
            if criterion != "technical_depth":
                await asyncio.sleep(0.01)
            return _mock_response(_valid_json(criterion))
        finally:
            state["active"] -= 1

    client = Mock()
    client.chat.completions.create.side_effect = create
    monkeypatch.setattr(
        "agents.scoring.scorer._get_async_groq_client", lambda key: client
    )

    result = asyncio.run(score_submission("sub-1", "text", groq_api_key="k"))

    assert isinstance(result, ScoringResult)
    assert result.submission_id == "sub-1"
    assert [s.criterion for s in result.scores] == CRITERIA_NAMES
    for s in result.scores:
        assert isinstance(s, CriterionScore)
        assert s.score == DISTINCT_SCORES[s.criterion]
        assert s.justification == f"{s.criterion} ok"


# --- Parallelism ---------------------------------------------------------------


def test_agents_run_concurrently_and_faster_than_sequential(monkeypatch):
    """All four calls overlap (max_active >= 2) and total wall-clock is
    meaningfully below the sequential sum — the roadmap's v0.5.0 DoD."""
    delay = 0.08  # seconds per fake call; sequential floor = 0.32s
    state = {"active": 0, "max_active": 0}

    async def create(**kwargs):
        system_content = kwargs["messages"][0]["content"]
        criterion = next(c for c in CRITERIA_NAMES if c in system_content)
        state["active"] += 1
        state["max_active"] = max(state["max_active"], state["active"])
        try:
            await asyncio.sleep(delay)
            return _mock_response(_valid_json(criterion))
        finally:
            state["active"] -= 1

    client = Mock()
    client.chat.completions.create.side_effect = create
    monkeypatch.setattr(
        "agents.scoring.scorer._get_async_groq_client", lambda key: client
    )

    start = time.perf_counter()
    result = asyncio.run(score_submission("sub-1", "text", groq_api_key="k"))
    elapsed = time.perf_counter() - start

    # Primary proof of parallelism: genuine overlap of in-flight calls.
    assert state["max_active"] >= 2
    # Secondary: wall-clock well under the sequential total (4 x delay).
    assert elapsed < len(CRITERIA_NAMES) * delay * 0.75
    assert [s.criterion for s in result.scores] == CRITERIA_NAMES


# --- Fail-closed semantics ------------------------------------------------------


def test_one_agent_failing_validation_fails_whole_run(monkeypatch):
    """No partial results: a persistently malformed agent aborts scoring."""
    client = _client_failing_criterion("innovation", "definitely not json")
    monkeypatch.setattr(
        "agents.scoring.scorer._get_async_groq_client", lambda key: client
    )

    with pytest.raises(RuntimeError, match="innovation"):
        asyncio.run(score_submission("sub-1", "text", groq_api_key="k"))


def test_persistent_rate_limit_propagates_as_rate_limit_error(monkeypatch):
    """An exhausted 429 surfaces as RateLimitError (a GroqError subclass),
    which the API route maps to 503."""
    rate_limit = RateLimitError("rate limited", response=Mock(), body=None)
    client = _client_failing_criterion("problem_fit", rate_limit)
    monkeypatch.setattr(
        "agents.scoring.scorer._get_async_groq_client", lambda key: client
    )

    with pytest.raises(RateLimitError):
        asyncio.run(score_submission("sub-1", "text", groq_api_key="k"))


def test_missing_groq_api_key_raises_value_error():
    with pytest.raises(ValueError, match="GROQ_API_KEY"):
        asyncio.run(score_submission("sub-1", "text"))


# --- Format independence (PPTX-only path) ---------------------------------------


def _pptx_bytes(text: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    box.text_frame.text = text
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_only_submission_scores_all_criteria(monkeypatch):
    """Roadmap DoD: a PPTX-only submission flows through real extraction
    into all four specialists without a PDF anywhere."""
    parsed = extract_text(_pptx_bytes(MARKER), "pptx")
    assert MARKER in parsed.raw_text

    seen_user_messages: list[str] = []

    async def create(**kwargs):
        system_content = kwargs["messages"][0]["content"]
        criterion = next(c for c in CRITERIA_NAMES if c in system_content)
        seen_user_messages.append(kwargs["messages"][1]["content"])
        return _mock_response(_valid_json(criterion))

    client = Mock()
    client.chat.completions.create.side_effect = create
    monkeypatch.setattr(
        "agents.scoring.scorer._get_async_groq_client", lambda key: client
    )

    scoring_text = build_scoring_text(parsed.raw_text, None)
    result = asyncio.run(score_submission("pptx-sub", scoring_text, groq_api_key="k"))

    assert [s.criterion for s in result.scores] == CRITERIA_NAMES
    assert all(s.justification for s in result.scores)
    # The extracted PPTX text reached every specialist's user message.
    assert len(seen_user_messages) == 4
    for message in seen_user_messages:
        assert MARKER in message


def test_image_descriptions_reach_every_agent_via_scoring_text(monkeypatch):
    """v0.3.5 merge behavior survives the split: descriptions ride along
    in the user message of ALL four agents."""
    descriptions = [
        {
            "page": 3,
            "classification": "architecture diagram",
            "description": "client talks to api over https",
        },
    ]
    captured: list[str] = []

    async def create(**kwargs):
        system_content = kwargs["messages"][0]["content"]
        criterion = next(c for c in CRITERIA_NAMES if c in system_content)
        captured.append(kwargs["messages"][1]["content"])
        return _mock_response(_valid_json(criterion))

    client = Mock()
    client.chat.completions.create.side_effect = create
    monkeypatch.setattr(
        "agents.scoring.scorer._get_async_groq_client", lambda key: client
    )

    scoring_text = build_scoring_text("plain text submission", descriptions)
    result = asyncio.run(score_submission("sub-img", scoring_text, groq_api_key="k"))

    assert [s.criterion for s in result.scores] == CRITERIA_NAMES
    assert len(captured) == 4
    for message in captured:
        assert "---IMAGE DESCRIPTIONS---" in message
