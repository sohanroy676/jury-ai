"""Tests for the submission upload endpoint."""

import io

import fitz
import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from backend.main import app
from backend.services import email as email_service
from backend.services import supabase

client = TestClient(app)

# Per-test registries backing the re-submission gate (v1.1.0):
# - _active_by_team maps lowercased team names to their ACTIVE row;
# - _insert_calls records every insert_submission invocation so tests
#   can assert the archive flag that reached the service layer.
# v1.2.0 adds _mailer_calls for confirmation-email dispatches.
_active_by_team: dict[str, dict] = {}
_insert_calls: list[dict] = []
_mailer_calls: list[dict] = []


@pytest.fixture(autouse=True)
def _mock_mailer(monkeypatch):
    """Route tests must never touch a real mail transport — the developer's
    .env may hold live SMTP credentials. Records dispatch kwargs; individual
    tests re-patch send_submission_confirmation for skip/fail outcomes."""
    _mailer_calls.clear()

    def fake_confirmation(**kwargs):
        _mailer_calls.append(kwargs)
        return email_service.EmailResult("sent", "sent")

    monkeypatch.setattr(
        email_service, "send_submission_confirmation", fake_confirmation
    )


@pytest.fixture(autouse=True)
def _mock_supabase(monkeypatch):
    """Mock the Supabase service layer so tests never hit the network."""

    def fake_upload(file_bytes, file_name, file_type):
        return f"https://example.supabase.co/storage/v1/object/public/submissions/{file_name}"

    insert_calls = _insert_calls
    insert_calls.clear()
    _active_by_team.clear()

    def fake_insert(
        team_name, file_url, file_type, team_email=None, supersedes_team=False, hackathon_id="default"
    ):
        insert_calls.append(
            {"team_email": team_email, "supersedes_team": supersedes_team}
        )
        return {
            "id": "00000000-0000-0000-0000-000000000001",
            "team_name": team_name,
            "file_url": file_url,
            "file_type": file_type,
            "status": "submitted",
            "team_email": team_email,
            "uploaded_at": "2026-08-26T10:00:00+00:00",
        }

    def fake_insert_parsed(
        submission_id,
        raw_text,
        sections,
        source_format,
        image_descriptions=None,
        hackathon_id="default",
    ):
        return {
            "id": "00000000-0000-0000-0000-000000000002",
            "submission_id": submission_id,
            "raw_text": raw_text,
            "sections": sections,
            "source_format": source_format,
            "image_descriptions": image_descriptions or [],
        }

    def fake_get_active(team_name):
        # Default: no active submission for any team.
        return _active_by_team.get(team_name.strip().lower())

    monkeypatch.setattr(supabase, "upload_submission_file", fake_upload)
    monkeypatch.setattr(supabase, "insert_submission", fake_insert)
    monkeypatch.setattr(supabase, "insert_parsed_submission", fake_insert_parsed)
    monkeypatch.setattr(supabase, "get_active_submission_by_team", fake_get_active)


def _pdf_bytes() -> bytes:
    """Build a real in-memory PDF with one page of text."""
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "JuryAI test proposal")
    data = doc.tobytes()
    doc.close()
    return data


def _pptx_bytes() -> bytes:
    """Build a real in-memory PPTX deck with one slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "JuryAI Test Deck"
    slide.placeholders[1].text = "Our solution"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_upload_pdf_success():
    """A valid PDF upload returns 201 and the stored row."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Alpha"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 201
    body = resp.json()
    assert body["team_name"] == "Team Alpha"
    assert body["file_type"] == "pdf"
    assert body["status"] == "submitted"
    assert body["file_url"].startswith("https://")


def test_upload_pptx_success():
    """A valid PPTX upload returns 201 and the stored row."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Beta"},
        files={"file": ("deck.pptx", _pptx_bytes(), "application/octet-stream")},
    )
    assert resp.status_code == 201
    body = resp.json()
    assert body["team_name"] == "Team Beta"
    assert body["file_type"] == "pptx"


def test_upload_rejects_unsupported_extension():
    """Uploading a .exe file is rejected with 400."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Gamma"},
        files={"file": ("malware.exe", b"MZ fake", "application/octet-stream")},
    )
    assert resp.status_code == 400
    assert "Unsupported file type" in resp.json()["detail"]


def test_upload_conflict_when_team_has_active_submission():
    """v1.1.0: a duplicate upload for a team with an ACTIVE submission is
    rejected with 409 unless the caller opts into replacing."""
    _active_by_team["team alpha"] = {
        "id": "00000000-0000-0000-0000-0000000000aa",
        "team_name": "Team Alpha",
        "uploaded_at": "2026-08-25T10:00:00+00:00",
    }

    # Exact-name duplicate...
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Alpha"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 409
    assert "already has an active submission" in resp.json()["detail"]

    # ...and the match is case-insensitive (identity = normalized name).
    resp_ci = client.post(
        "/api/submissions",
        data={"team_name": "TEAM ALPHA"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp_ci.status_code == 409


def test_upload_replace_existing_archives_previous():
    """v1.1.0: replace_existing=true lets the re-upload through and flags
    the insert so the service layer archives the superseded row."""
    _active_by_team["team alpha"] = {
        "id": "00000000-0000-0000-0000-0000000000aa",
        "team_name": "Team Alpha",
        "uploaded_at": "2026-08-25T10:00:00+00:00",
    }

    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Alpha", "replace_existing": "true"},
        files={"file": ("proposal-v2.pdf", _pdf_bytes(), "application/pdf")},
    )

    assert resp.status_code == 201, resp.text
    assert len(_insert_calls) == 1
    assert _insert_calls[0]["supersedes_team"] is True


def test_upload_no_conflict_for_fresh_team():
    """A brand-new team uploads normally; no archive flag is sent."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Brand New Team"},
        files={"file": ("first.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 201
    assert _insert_calls[-1]["supersedes_team"] is False


def test_upload_rejects_docx():
    """Uploading a .docx file is rejected with 400."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Delta"},
        files={"file": ("notes.docx", b"PK fake docx", "application/octet-stream")},
    )
    assert resp.status_code == 400


def test_upload_rejects_empty_team_name():
    """A blank team name is rejected with 400."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "   "},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 400
    assert "Team name" in resp.json()["detail"]


def test_upload_rejects_empty_file():
    """An empty file is rejected with 400."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Epsilon"},
        files={"file": ("empty.pdf", b"", "application/pdf")},
    )
    assert resp.status_code == 400
    assert "empty" in resp.json()["detail"].lower()


def test_upload_rejects_oversized_file(monkeypatch):
    """A file larger than MAX_UPLOAD_BYTES is rejected with 413."""
    from backend.routes import submissions as routes

    monkeypatch.setattr(routes, "MAX_UPLOAD_BYTES", 10)  # 10 bytes
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Zeta"},
        files={"file": ("big.pdf", b"x" * 100, "application/pdf")},
    )
    assert resp.status_code == 413
    assert "too large" in resp.json()["detail"].lower()


def test_upload_rejects_corrupt_pdf():
    """A file with a .pdf extension but unparseable content is rejected with 422."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Corrupt"},
        files={"file": ("broken.pdf", b"this is not a real pdf", "application/pdf")},
    )
    assert resp.status_code == 422
    assert "Could not parse file" in resp.json()["detail"]


# --- Confirmation notifications (v1.2.0) -----------------------------------------


def test_upload_sends_confirmation_email():
    """A valid upload with an address dispatches exactly one confirmation."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Alpha", "team_email": "alpha@example.com"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 201
    body = resp.json()
    assert body["team_email"] == "alpha@example.com"
    assert body["notification"]["confirmation_email"]["status"] == "sent"

    assert len(_mailer_calls) == 1
    assert _mailer_calls[0]["team_email"] == "alpha@example.com"
    assert _mailer_calls[0]["team_name"] == "Team Alpha"
    assert _mailer_calls[0]["file_type"] == "pdf"


def test_upload_without_email_passes_blank_address_to_mailer():
    """Blank input is optional at the API level and travels as ''."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "No Mail"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 201
    assert len(_mailer_calls) == 1
    assert _mailer_calls[0]["team_email"] == ""


def test_upload_surfaces_skipped_notification(monkeypatch):
    """The response reports the service's skip outcome verbatim."""

    def skipping(**kwargs):
        return email_service.EmailResult("skipped", "no_valid_recipient", "no addr")

    monkeypatch.setattr(email_service, "send_submission_confirmation", skipping)

    resp = client.post(
        "/api/submissions",
        data={"team_name": "Legacy Team"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 201
    notif = resp.json()["notification"]["confirmation_email"]
    assert notif["status"] == "skipped"
    assert notif["reason"] == "no_valid_recipient"


def test_upload_rejects_malformed_email_before_persisting():
    """A malformed address is a 400 BEFORE any row or mail exists."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Typo Team", "team_email": "not-an-email"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 400
    assert "email" in resp.json()["detail"].lower()
    assert _insert_calls == []
    assert _mailer_calls == []


def test_parse_failure_never_emails():
    """A corrupt deck must not produce a confirmation (422 path)."""
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Corrupt Team"},
        files={"file": ("broken.pdf", b"not a real pdf", "application/pdf")},
    )
    assert resp.status_code == 422
    assert _mailer_calls == []


def test_conflict_gate_never_emails():
    """The 409 duplicate gate fires long before any notification."""
    _active_by_team["team alpha"] = {
        "id": "00000000-0000-0000-0000-0000000000aa",
        "team_name": "Team Alpha",
        "uploaded_at": "2026-08-25T10:00:00+00:00",
    }
    resp = client.post(
        "/api/submissions",
        data={"team_name": "Team Alpha"},
        files={"file": ("proposal.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 409
    assert _mailer_calls == []


def test_replace_flow_confirms_the_new_active_row():
    """Re-submission archives the old row AND confirms the new one."""
    _active_by_team["team alpha"] = {
        "id": "00000000-0000-0000-0000-0000000000aa",
        "team_name": "Team Alpha",
        "uploaded_at": "2026-08-25T10:00:00+00:00",
    }
    resp = client.post(
        "/api/submissions",
        data={
            "team_name": "Team Alpha",
            "replace_existing": "true",
            "team_email": "alpha-new@example.com",
        },
        files={"file": ("proposal-v2.pdf", _pdf_bytes(), "application/pdf")},
    )
    assert resp.status_code == 201
    assert _insert_calls[0]["supersedes_team"] is True
    assert _insert_calls[0]["team_email"] == "alpha-new@example.com"
    assert len(_mailer_calls) == 1
    assert _mailer_calls[0]["team_email"] == "alpha-new@example.com"


def test_health():
    """The health endpoint returns ok."""
    resp = client.get("/health")
    assert resp.status_code == 200
    assert resp.json() == {"status": "ok"}


def test_cors_allows_frontend_origin():
    """The backend must return CORS headers so the browser lets the
    frontend read the response (regression test for the
    'Network error — could not reach the backend' bug)."""
    resp = client.options(
        "/api/submissions",
        headers={
            "Origin": "http://localhost:3000",
            "Access-Control-Request-Method": "POST",
            "Access-Control-Request-Headers": "Content-Type",
        },
    )
    assert resp.headers.get("access-control-allow-origin") == "http://localhost:3000"
    assert "POST" in resp.headers.get("access-control-allow-methods", "")
    assert (
        "content-type" in resp.headers.get("access-control-allow-headers", "").lower()
    )


def test_cors_preflight_allows_put_for_rubrics():
    """Saving rubric weights from the dashboard sends PUT with a JSON
    body, so browsers preflight it. Starlette's CORSMiddleware answers
    a preflight whose requested method is not whitelisted with
    400 'Disallowed CORS method' (surfaced in the UI as 'Network error
    could not reach the backend'). Regression test for the v0.6.0
    PUT /api/rubrics route being unreachable while CORS allowed only
    GET/POST."""
    resp = client.options(
        "/api/rubrics/default",
        headers={
            "Origin": "http://localhost:3000",
            "Access-Control-Request-Method": "PUT",
            "Access-Control-Request-Headers": "Content-Type",
        },
    )
    assert resp.status_code == 200
    assert resp.headers.get("access-control-allow-origin") == "http://localhost:3000"
    assert "PUT" in resp.headers.get("access-control-allow-methods", "")
    assert (
        "content-type" in resp.headers.get("access-control-allow-headers", "").lower()
    )


def test_cors_rejects_unlisted_origin():
    """An origin not in FRONTEND_URL should not get an allow-origin header."""
    resp = client.options(
        "/api/submissions",
        headers={"Origin": "http://evil.example.com"},
    )
    assert resp.headers.get("access-control-allow-origin") is None
