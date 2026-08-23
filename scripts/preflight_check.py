"""Pre-flight check for v0.3.5 end-to-end verification.

1. Reports which required env vars are SET/MISSING (never prints values).
2. Inspects a submission file's composition: slides/pages, text volume,
   content-area images (via the real extractor) and master/layout
   images (which the structural filter must exclude).

Run: python scripts/preflight_check.py <file>
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from dotenv import load_dotenv

load_dotenv(os.path.join(os.path.dirname(os.path.dirname(__file__)), ".env"))


def check_env() -> bool:
    print("=== Environment variables ===")
    required = [
        "SUPABASE_URL",
        "SUPABASE_ANON_KEY",
        "SUPABASE_SERVICE_ROLE_KEY",
        "GROQ_API_KEY",
    ]
    ok = True
    for name in required:
        present = bool(os.getenv(name, "").strip())
        print(f"  {name}: {'SET' if present else 'MISSING'}")
        ok = ok and present
    return ok


def inspect_file(path: str) -> None:
    print(f"\n=== File inspection: {path} ===")
    if not os.path.isfile(path):
        print("  FILE NOT FOUND")
        raise SystemExit(1)

    size_mb = os.path.getsize(path) / (1024 * 1024)
    print(f"  Size: {size_mb:.2f} MB")

    ext = os.path.splitext(path)[1].lower()
    source_format = {".pdf": "pdf", ".pptx": "pptx"}.get(ext)
    if source_format is None:
        print(f"  Unsupported extension '{ext}' - only .pdf/.pptx")
        raise SystemExit(1)

    with open(path, "rb") as f:
        file_bytes = f.read()

    # Text side.
    from agents.parsing.extractor import extract_text

    parsed = extract_text(file_bytes, source_format)
    total_text = sum(len(s["text"]) for s in parsed.sections)
    print(f"  Format: {source_format}")
    print(f"  Sections (slides/pages): {len(parsed.sections)}")
    print(f"  Total extracted text chars: {total_text}")

    # Image side - content areas (what WILL be processed).
    from agents.parsing.images.dedupe import compute_phash
    from agents.parsing.images.extract import extract_images

    images = extract_images(file_bytes, source_format)
    print(f"  Content-area images found: {len(images)}")
    by_page: dict[int, int] = {}
    unhashable = 0
    hashes: list[str] = []
    for img in images:
        by_page[img.page_number] = by_page.get(img.page_number, 0) + 1
        ph = compute_phash(img.image_bytes)
        if ph is None:
            unhashable += 1
        else:
            hashes.append(ph)
    print(f"    per page/slide: {dict(sorted(by_page.items()))}")
    print(f"    undecodable (will be skipped): {unhashable}")
    unique = len(set(hashes))
    dupes = len(hashes) - unique
    print(
        f"    unique pHashes after dedupe: {unique}"
        f" ({dupes} within-submission near-duplicates will be dropped)"
    )

    # Master/layout images (what must be EXCLUDED).
    if source_format == "pptx":
        import io

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(file_bytes))

        def count_pictures(shapes) -> int:
            n = 0
            for shape in shapes:
                try:
                    st = shape.shape_type
                except Exception as exc:  # noqa: BLE001 - inspection only; skip unreadable shapes
                    print(f"    (skipping shape with unreadable type: {exc})")
                    continue
                if st == MSO_SHAPE_TYPE.GROUP:
                    n += count_pictures(shape.shapes)
                elif st == MSO_SHAPE_TYPE.PICTURE:
                    n += 1
            return n

        master_imgs = sum(count_pictures(m.shapes) for m in prs.slide_masters)
        layout_imgs = sum(
            count_pictures(layout.shapes)
            for m in prs.slide_masters
            for layout in m.slide_layouts
        )
        print(
            f"  Master/layout images (must be EXCLUDED by structural"
            f" filter): master={master_imgs}, layouts={layout_imgs}"
        )


def main() -> None:
    if len(sys.argv) != 2:
        raise SystemExit("Usage: python scripts/preflight_check.py <file>")

    env_ok = check_env()
    inspect_file(sys.argv[1])

    print("\n=== Pre-flight verdict ===")
    print("Env OK" if env_ok else "ENV INCOMPLETE - fill .env before running")


if __name__ == "__main__":
    main()
